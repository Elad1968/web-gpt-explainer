import pptx
import pptx.presentation
import pptx.table


class PresentationHelper:
    @staticmethod
    def extract_text_from_table(table: pptx.table.Table) -> str:
        table_text: list[list[str]] = []
        for row in table.rows:
            row_text: list[str] = []
            for cell in row.cells:
                cell_text: str = cell.text.strip()
                row_text.append(cell_text)
            table_text.append(row_text)
        text: str = ""
        for row in table_text:
            text += '\t'.join(row)
            text += '\n'
        return text

    @staticmethod
    def get_all_text_from_presentation(presentation: pptx.presentation.Presentation) -> list[str]:
        result: list[str] = []
        for slide in presentation.slides:
            text: str = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text.strip()
                if text and not text[-1].isspace():
                    text += ' '
                if shape.has_table:
                    text += PresentationHelper.extract_text_from_table(shape.table)
                if text and not text[-1].isspace():
                    text += ' '
            result.append(text)
        return [slide_text for slide_text in result if not slide_text.isspace() and slide_text]

    @staticmethod
    def get_presentation(path: str) -> pptx.presentation.Presentation:
        return pptx.Presentation(path)
